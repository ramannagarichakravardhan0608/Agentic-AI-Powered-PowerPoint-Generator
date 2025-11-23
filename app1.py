
from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Use slide layout 1 (Title and Content) for all slides
title_content_layout = prs.slide_layouts[1]

# Slide 1: What is Artificial Intelligence?
slide1 = prs.slides.add_slide(title_content_layout)
title1 = slide1.shapes.title
title1.text = "What is Artificial Intelligence?"
body1 = slide1.shapes.placeholders[1]
tf1 = body1.text_frame
tf1.clear() # Clear any default text
p1 = tf1.add_paragraph()
p1.text = "Simulation of human intelligence in machines."
p1 = tf1.add_paragraph()
p1.text = "Enables machines to learn, reason, perceive, and act."
p1 = tf1.add_paragraph()
p1.text = "A broad field encompassing Machine Learning, Deep Learning, NLP, etc."
p1 = tf1.add_paragraph()
p1.text = "Aims to create systems that can perform tasks requiring human intellect."

# Slide 2: What is n8n?
slide2 = prs.slides.add_slide(title_content_layout)
title2 = slide2.shapes.title
title2.text = "What is n8n?"
body2 = slide2.shapes.placeholders[1]
tf2 = body2.text_frame
tf2.clear()
p2 = tf2.add_paragraph()
p2.text = "An open-source workflow automation tool."
p2 = tf2.add_paragraph()
p2.text = "Helps connect applications and services to automate tasks."
p2 = tf2.add_paragraph()
p2.text = "Features a visual workflow designer with a node-based interface."
p2 = tf2.add_paragraph()
p2.text = "Supports various integrations, including AI and generative services."
p2 = tf2.add_paragraph()
p2.text = "Enables creation of complex automations without extensive coding."

# Slide 3: AI vs Generative AI – Key Differences
slide3 = prs.slides.add_slide(title_content_layout)
title3 = slide3.shapes.title
title3.text = "AI vs Generative AI – Key Differences"
body3 = slide3.shapes.placeholders[1]
tf3 = body3.text_frame
tf3.clear()
p3 = tf3.add_paragraph()
p3.text = "Artificial Intelligence (AI):"
p3.level = 0
p3 = tf3.add_paragraph()
p3.text = "Broad concept of machines simulating human intelligence."
p3.level = 1
p3 = tf3.add_paragraph()
p3.text = "Focuses on problem-solving, decision-making, pattern recognition."
p3.level = 1
p3 = tf3.add_paragraph()
p3.text = "Generative AI:"
p3.level = 0
p3 = tf3.add_paragraph()
p3.text = "A sub-field of AI focused on creating new content (text, images, audio)."
p3.level = 1
p3 = tf3.add_paragraph()
p3.text = "Utilizes models like LLMs and diffusion models to generate novel outputs."
p3.level = 1

# Slide 4: Types of Artificial Intelligence
slide4 = prs.slides.add_slide(title_content_layout)
title4 = slide4.shapes.title
title4.text = "Types of Artificial Intelligence"
body4 = slide4.shapes.placeholders[1]
tf4 = body4.text_frame
tf4.clear()
p4 = tf4.add_paragraph()
p4.text = "Narrow AI (Weak AI):"
p4.level = 0
p4 = tf4.add_paragraph()
p4.text = "Designed for specific tasks (e.g., Siri, recommendation systems)."
p4.level = 1
p4 = tf4.add_paragraph()
p4.text = "General AI (Strong AI):"
p4.level = 0
p4 = tf4.add_paragraph()
p4.text = "Hypothetical AI with human-level cognitive abilities across tasks."
p4.level = 1
p4 = tf4.add_paragraph()
p4.text = "Super AI:"
p4.level = 0
p4 = tf4.add_paragraph()
p4.text = "Hypothetical AI surpassing human intelligence in all aspects."
p4.level = 1

# Slide 5: Applications of AI
slide5 = prs.slides.add_slide(title_content_layout)
title5 = slide5.shapes.title
title5.text = "Applications of AI"
body5 = slide5.shapes.placeholders[1]
tf5 = body5.text_frame
tf5.clear()
p5 = tf5.add_paragraph()
p5.text = "Healthcare: Disease diagnosis, drug discovery, personalized medicine."
p5 = tf5.add_paragraph()
p5.text = "Finance: Fraud detection, algorithmic trading, risk assessment."
p5 = tf5.add_paragraph()
p5.text = "Retail: Recommendation systems, inventory management, chatbots."
p5 = tf5.add_paragraph()
p5.text = "Automotive: Self-driving cars, predictive maintenance."
p5 = tf5.add_paragraph()
p5.text = "Education: Personalized learning paths, intelligent tutoring systems."
p5 = tf5.add_paragraph()
p5.text = "Customer Service: Virtual assistants, sentiment analysis."

# Slide 6: Future of Artificial Intelligence
slide6 = prs.slides.add_slide(title_content_layout)
title6 = slide6.shapes.title
title6.text = "Future of Artificial Intelligence"
body6 = slide6.shapes.placeholders[1]
tf6 = body6.text_frame
tf6.clear()
p6 = tf6.add_paragraph()
p6.text = "Ethical AI: Focus on fairness, transparency, and accountability."
p6 = tf6.add_paragraph()
p6.text = "Agentic AI: Development of autonomous, goal-oriented AI systems."
p6 = tf6.add_paragraph()
p6.text = "Human-AI Collaboration: Enhanced partnership for complex problem-solving."
p6 = tf6.add_paragraph()
p6.text = "Specialized AI: More sophisticated and narrowly focused solutions."
p6 = tf6.add_paragraph()
p6.text = "Regulation: Increasing governmental and societal oversight."
p6 = tf6.add_paragraph()
p6.text = "Quantum AI: Integration with quantum computing for unprecedented power."

# Slide 7: Conclusion
slide7 = prs.slides.add_slide(title_content_layout)
title7 = slide7.shapes.title
title7.text = "Conclusion"
body7 = slide7.shapes.placeholders[1]
tf7 = body7.text_frame
tf7.clear()
p7 = tf7.add_paragraph()
p7.text = "AI is a transformative technology reshaping industries and daily life."
p7 = tf7.add_paragraph()
p7.text = "From basic automation to generative capabilities, its impact is growing."
p7 = tf7.add_paragraph()
p7.text = "The future holds promise for more autonomous and intelligent systems (Agentic AI)."
p7 = tf7.add_paragraph()
p7.text = "Ethical considerations and responsible development are crucial."
p7 = tf7.add_paragraph()
p7.text = "Continuous learning and adaptation will define our interaction with AI."

prs.save("presentation.pptx")
